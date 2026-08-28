"""shortform_template.json 구성에 맞춰 영상 제작용 PPTX를 생성한다.
각 슬라이드 = 하나의 씬. 배경 이미지는 렌더링된 프레임 PNG를 그대로 사용하고,
스피커 노트에는 해당 구간 대본을 넣어 촬영/더빙 시 참고할 수 있게 한다.
"""
import json
import os

from pptx import Presentation
from pptx.util import Inches

from common import TEMPLATES_DIR

SLIDE_W_IN = 6.08  # 1080/1920 비율 유지 (9:16)
SLIDE_H_IN = 10.8


def build_pptx(content, frame_paths, out_path):
    with open(os.path.join(TEMPLATES_DIR, "shortform_template.json"), encoding="utf-8") as f:
        template = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]

    script = content["script"]
    scenes_by_id = {s["id"]: s for s in template["scenes"]}
    frame_by_scene = {}
    for path in frame_paths:
        name = os.path.basename(path)  # frame_01_hook.png
        scene_id = name.split("_", 2)[2].rsplit(".", 1)[0]
        frame_by_scene[scene_id] = path

    note_map = {
        "hook": script.get("hook", ""),
        "paper_intro": script.get("intro", ""),
        "key_finding_1": (script.get("highlights") or [""])[0],
        "key_finding_2": (script.get("highlights") or ["", ""])[1] if len(script.get("highlights", [])) > 1 else "",
        "my_take": "본인 코멘트를 자유롭게 추가하세요.",
        "cta": script.get("cta", ""),
    }

    for scene in sorted(template["scenes"], key=lambda s: s["order"]):
        img_path = frame_by_scene.get(scene["id"])
        if not img_path or not os.path.exists(img_path):
            continue
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        notes = slide.notes_slide.notes_text_frame
        notes.text = (
            f"[{scene['id']}] 약 {scene['duration_sec']}초\n"
            f"대본: {note_map.get(scene['id'], '')}\n"
            f"연출 노트: {scene.get('notes', '')}"
        )

    # 참고 자료 슬라이드 (원문 출처)
    ref_slide = prs.slides.add_slide(blank_layout)
    tb = ref_slide.shapes.add_textbox(Inches(0.4), Inches(0.4), prs.slide_width - Inches(0.8), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = "참고 자료 출처"
    p = tf.add_paragraph()
    p.text = content["topic"].get("title", "")
    p2 = tf.add_paragraph()
    p2.text = content["topic"].get("url", "")

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path
